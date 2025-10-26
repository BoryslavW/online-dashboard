from pptx import Presentation
from pptx.util import Inches
from io import BytesIO  # This is for PPT and charts—keep it
import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
import io  # This is for StringIO in CSV download—MOVE IT HERE if it's elsewhere
import base64  # For chart downloads
import time    # For progress bar

# Define the presentation creation function (outside the main block)
def create_presentation(df, sector_counts, region_counts, size_dist, fig1, fig2, fig3):
    prs = Presentation()
    
    # Title Slide
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Company Data Analysis Report"
    subtitle.text = f"Generated on {pd.Timestamp.now().strftime('%Y-%m-%d %H:%M')}\nTotal Companies: {len(df)}"
    
    # Summary Stats Slide
    slide_layout = prs.slide_layouts[1]  # Title and content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Summary Statistics"
    content = slide.placeholders[1]
    content.text = df.describe(include='all').to_string()  # Simple text summary
    
    # Sector Bar Chart Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Top 10 Sectors by Company Count"
    img_buf = BytesIO()
    fig1.savefig(img_buf, format='png', bbox_inches='tight')
    img_buf.seek(0)
    left = Inches(1)
    top = Inches(2)
    slide.shapes.add_picture(img_buf, left, top, height=Inches(4))
    
    # Region Pie Chart Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Top 5 Countries Distribution"
    img_buf = BytesIO()
    fig2.savefig(img_buf, format='png', bbox_inches='tight')
    img_buf.seek(0)
    slide.shapes.add_picture(img_buf, left, top, height=Inches(4))
    
    # Size Histogram Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Company Size Category Distribution"
    img_buf = BytesIO()
    fig3.savefig(img_buf, format='png', bbox_inches='tight')
    img_buf.seek(0)
    slide.shapes.add_picture(img_buf, left, top, height=Inches(4))
    
    # Save to buffer for download
    pptx_buf = BytesIO()
    prs.save(pptx_buf)
    pptx_buf.seek(0)
    return pptx_buf

# Streamlit app title and instructions
st.title("Company Data Analyzer")
st.markdown("**Welcome!** Upload a CSV file with columns: `company_name`, `focus`, `sector`, `region`, and `size` to analyze company data. The `size` column should contain values like '1-10', '11-50', etc., or 'Unknown'. Click 'Browse files' to start!")

# File uploader with format warning
uploaded_file = st.file_uploader("Choose a CSV file", type="csv")
if uploaded_file is not None:
    # Read the uploaded CSV with progress bar
    progress_bar = st.progress(0)
    for i in range(100):
        time.sleep(0.01)  # Simulate work (adjust for larger files)
        progress_bar.progress(i + 1)
    progress_bar.empty()  # Hide when done
    
    df = pd.read_csv(uploaded_file)
    required_columns = ['company_name', 'focus', 'sector', 'region', 'size']
    if not all(col in df.columns for col in required_columns):
        st.error("Error: Your CSV must contain all columns: `company_name`, `focus`, `sector`, `region`, `size`. Please check your file and try again.")
        st.stop()
    
    # Your cleaning logic
    df = df.drop_duplicates(subset=['company_name'])
    df.fillna({
        'focus': 'Unknown',
        'sector': 'Unknown',
        'region': 'Unknown',
        'size': 'Unknown'
    }, inplace=True)
    df = df.dropna(how='all', subset=['sector', 'region', 'size'])
    for col in ['focus', 'sector', 'region', 'size']:
        df[col] = df[col].astype(str).str.lower().str.strip()
    valid_sizes = ['1-10', '11-50', '51-200', '201-500', '501-1k', 'unknown']
    df['size'] = df['size'].apply(lambda x: x if x in valid_sizes else 'unknown')
    def categorize_size(size):
        if size in ['1-10', '11-50']:
            return 'Small'
        elif size in ['51-200', '201-500']:
            return 'Medium'
        elif size in ['501-1k']:
            return 'Large'
        else:
            return 'Unknown'
    df['size_category'] = df['size'].apply(categorize_size)
    df['country'] = df['region'].apply(lambda x: x.split(',')[-1].strip() if ',' in x else x)
    
    # Interactive sector filter
    st.subheader("Filter by Sector")
    unique_sectors = df['sector'].dropna().unique().tolist()
    selected_sector = st.selectbox("Choose a sector to filter (or 'All')", ['All'] + unique_sectors)
    if selected_sector != 'All':
        df = df[df['sector'] == selected_sector]
    
    # Analysis
    sector_counts = df['sector'].value_counts().reset_index(name='count').head(10)
    region_counts = df['region'].value_counts().reset_index(name='count').head(10)
    size_dist = df['size_category'].value_counts().reset_index(name='count')
    
    # Display summary stats
    st.subheader("Summary Stats")
    st.dataframe(df.describe(include='all'))
    
    # Display tables
    st.subheader("Top 10 Sectors by Company Count")
    st.dataframe(sector_counts)
    st.subheader("Region Counts")
    st.dataframe(region_counts)
    
    # Generate and display plots
    st.subheader("Visualizations")
    # Top 10 Sectors Bar Plot
    fig1, ax1 = plt.subplots(figsize=(10, 6))
    sns.barplot(x='count', y='sector', data=sector_counts, ax=ax1)
    ax1.set_title('Top 10 Sectors by Company Count')
    ax1.set_xlabel('Count')
    ax1.set_ylabel('Sector')
    plt.tight_layout()
    st.pyplot(fig1)
    # Download button for sector chart
    def get_image_download_link(fig, filename):
        buf = BytesIO()
        fig.savefig(buf, format="png", bbox_inches='tight')
        buf.seek(0)
        b64 = base64.b64encode(buf.read()).decode()
        href = f'<a href="data:file/png;base64,{b64}" download="{filename}.png">Download {filename} Chart</a>'
        return st.markdown(href, unsafe_allow_html=True)
    get_image_download_link(fig1, "sectors_bar")
    
    # Top 5 Countries Pie Chart
    country_counts = df['country'].value_counts().head(5)
    fig2, ax2 = plt.subplots(figsize=(8, 8))
    ax2.pie(country_counts, labels=country_counts.index, autopct='%1.1f%%', startangle=90)
    ax2.set_title('Top 5 Countries Distribution')
    plt.tight_layout()
    st.pyplot(fig2)
    get_image_download_link(fig2, "regions_pie")
    
    # Size Category Histogram
    fig3, ax3 = plt.subplots(figsize=(8, 6))
    sns.histplot(df['size_category'], kde=False, ax=ax3)
    ax3.set_title('Company Size Category Distribution')
    ax3.set_xlabel('Size Category')
    ax3.set_ylabel('Count')
    plt.tight_layout()
    st.pyplot(fig3)
    get_image_download_link(fig3, "sizes_histogram")
    
    # Save and offer download for cleaned CSV
    cleaned_csv = io.StringIO()
    df.to_csv(cleaned_csv, index=False)
    st.download_button(
        label="Download Cleaned CSV",
        data=cleaned_csv.getvalue(),
        file_name="cleaned_company_data.csv",
        mime="text/csv"
    )
    
    # Generate and download PowerPoint
    if st.button("Generate PowerPoint Presentation"):
        pptx_buffer = create_presentation(df, sector_counts, region_counts, size_dist, fig1, fig2, fig3)
        st.download_button(
            label="Download Presentation (.pptx)",
            data=pptx_buffer.getvalue(),
            file_name="company_analysis.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
else:
    st.info("Please upload a CSV file to analyze.")

